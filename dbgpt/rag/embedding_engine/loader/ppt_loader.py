from typing import List, Optional

from langchain.docstore.document import Document
from langchain.document_loaders.base import BaseLoader
from pptx import Presentation


class PPTLoader(BaseLoader):
    """Load PPT files."""

    def __init__(self, file_path: str, encoding: Optional[str] = None):
        """Initialize with file path."""
        self.file_path = file_path
        self.encoding = encoding

    def load(self) -> List[Document]:
        """Load from file path."""
        pr = Presentation(self.file_path)
        docs = []
        for slide in pr.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    docs.append(
                        Document(
                            page_content=shape.text, metadata={"source": slide.slide_id}
                        )
                    )
        return docs
