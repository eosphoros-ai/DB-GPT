"""PPTX Knowledge."""

from typing import Any, Dict, List, Optional, Union

from dbgpt.core import Document
from dbgpt.rag.knowledge.base import (
    ChunkStrategy,
    DocumentType,
    Knowledge,
    KnowledgeType,
)


class PPTXKnowledge(Knowledge):
    """PPTX Knowledge."""

    def __init__(
        self,
        file_path: Optional[str] = None,
        knowledge_type: KnowledgeType = KnowledgeType.DOCUMENT,
        loader: Optional[Any] = None,
        language: Optional[str] = "zh",
        metadata: Optional[Dict[str, Union[str, List[str]]]] = None,
        **kwargs: Any,
    ) -> None:
        """Create PPTX knowledge with PDF Knowledge arguments.

        Args:
            file_path:(Optional[str]) file path
            knowledge_type:(KnowledgeType) knowledge type
            loader:(Optional[Any]) loader
        """
        super().__init__(
            path=file_path,
            knowledge_type=knowledge_type,
            data_loader=loader,
            metadata=metadata,
            **kwargs,
        )
        self._language = language

    def _load(self) -> List[Document]:
        """Load pdf document from loader."""
        if self._loader:
            documents = self._loader.load()
        else:
            from pptx import Presentation

            pr = Presentation(self._path)
            docs = []
            for slide in pr.slides:
                content = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        content += shape.text
                metadata = {"source": self._path}
                if self._metadata:
                    metadata.update(self._metadata)  # type: ignore
                docs.append(Document(content=content, metadata=metadata))
            return docs
        return [Document.langchain2doc(lc_document) for lc_document in documents]

    @classmethod
    def support_chunk_strategy(cls) -> List[ChunkStrategy]:
        """Return support chunk strategy.

        Returns:
            List[ChunkStrategy]: support chunk strategy
        """
        return [
            ChunkStrategy.CHUNK_BY_SIZE,
            ChunkStrategy.CHUNK_BY_PAGE,
            ChunkStrategy.CHUNK_BY_SEPARATOR,
        ]

    @classmethod
    def default_chunk_strategy(cls) -> ChunkStrategy:
        """Return default chunk strategy.

        Returns:
            ChunkStrategy: default chunk strategy
        """
        return ChunkStrategy.CHUNK_BY_SIZE

    @classmethod
    def type(cls) -> KnowledgeType:
        """Knowledge type of PPTX.

        Returns:
            KnowledgeType: knowledge type
        """
        return KnowledgeType.DOCUMENT

    @classmethod
    def document_type(cls) -> DocumentType:
        """Document type of PPTX.

        Returns:
            DocumentType: document type
        """
        return DocumentType.PPTX
